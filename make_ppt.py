from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BLACK = RGBColor(0x11, 0x11, 0x11)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_TEXT = RGBColor(0x44, 0x44, 0x44)
GRAY_SUB = RGBColor(0x66, 0x66, 0x66)
GRAY_LIGHT = RGBColor(0x88, 0x88, 0x88)
GRAY_BG = RGBColor(0xF5, 0xF5, 0xF5)
GRAY_BORDER = RGBColor(0xE0, 0xE0, 0xE0)
BLUE_FK = RGBColor(0x00, 0x66, 0xCC)

FONT = 'Pretendard Std'
FONT_FALLBACK = 'Malgun Gothic'

def get_font():
    return FONT

def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name=None):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name or get_font()
    p.alignment = alignment
    return txBox

def add_paragraph(text_frame, text, font_size=14, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, space_before=0, space_after=0):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = get_font()
    p.alignment = alignment
    if space_before:
        p.space_before = Pt(space_before)
    if space_after:
        p.space_after = Pt(space_after)
    return p

def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(height))
    shape.fill.background() if fill_color is None else None
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    # Smaller corner radius
    shape.adjustments[0] = 0.02
    return shape

def add_header_bar(slide, section_title):
    """Add the black header bar with section title and logo"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Emu(438912))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLACK
    bar.line.fill.background()

    # Section title left
    add_textbox(slide, Emu(438912), Emu(100000), Emu(4000000), Emu(338912),
                section_title, font_size=11, color=RGBColor(0xAA, 0xAA, 0xAA))

    # Logo right
    add_textbox(slide, Emu(10800000), Emu(100000), Emu(1800000), Emu(338912),
                'aS  Always Solving', font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.RIGHT)

def add_slide_number(slide, num):
    add_textbox(slide, Emu(11800000), Emu(6500000), Emu(500000), Emu(300000),
                str(num), font_size=10, color=GRAY_LIGHT, alignment=PP_ALIGN.RIGHT)

def add_title_subtitle(slide, title, subtitle, top_offset=600000):
    add_textbox(slide, Emu(585216), Emu(top_offset), Emu(10000000), Emu(500000),
                title, font_size=28, bold=True, color=BLACK)
    if subtitle:
        add_textbox(slide, Emu(585216), Emu(top_offset + 420000), Emu(10000000), Emu(300000),
                    subtitle, font_size=14, color=GRAY_SUB)

# ═══════════════════════════════════════
# SLIDE 1: Cover
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# Accent line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(585216), Emu(2200000), Emu(548640), Emu(36576))
line.fill.solid()
line.fill.fore_color.rgb = BLACK
line.line.fill.background()

# Title
add_textbox(slide, Emu(585216), Emu(2400000), Emu(10000000), Emu(900000),
            '미국 정부 입찰 정보\n자동 수집 시스템', font_size=36, bold=True, color=BLACK)

# Subtitle
add_textbox(slide, Emu(585216), Emu(3500000), Emu(5000000), Emu(300000),
            '프로젝트 제안서', font_size=18, color=GRAY_SUB)

# Meta info
meta_top = Emu(4200000)
add_textbox(slide, Emu(585216), meta_top, Emu(2500000), Emu(250000),
            '제안사  asolve (에이솔브)', font_size=12, color=GRAY_LIGHT)
add_textbox(slide, Emu(3500000), meta_top, Emu(2500000), Emu(250000),
            '발표일  2026년 04월 06일', font_size=12, color=GRAY_LIGHT)
add_textbox(slide, Emu(6200000), meta_top, Emu(2500000), Emu(250000),
            '담당자  유영주', font_size=12, color=GRAY_LIGHT)


# ═══════════════════════════════════════
# SLIDE 2: Introduction
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, '1. 소개')
add_slide_number(slide, 2)

add_textbox(slide, Emu(585216), Emu(600000), Emu(5000000), Emu(400000),
            '유영주', font_size=28, bold=True)

# Label badge
badge = add_rect(slide, Emu(585216), Emu(1050000), Emu(2800000), Emu(280000), fill_color=GRAY_BG, border_color=GRAY_BORDER)
add_textbox(slide, Emu(650000), Emu(1060000), Emu(2700000), Emu(260000),
            '풀스택 개발자  ·  21년 경력', font_size=11, bold=True, color=GRAY_SUB)

# Left column - Career Summary
left_x = Emu(585216)
col_top = Emu(1500000)

add_textbox(slide, left_x, col_top, Emu(5000000), Emu(300000),
            '경력 요약', font_size=16, bold=True)

# Divider
div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, Emu(1800000), Emu(365760), Emu(27432))
div.fill.solid()
div.fill.fore_color.rgb = BLACK
div.line.fill.background()

txBox = slide.shapes.add_textbox(left_x, Emu(1900000), Emu(5400000), Emu(2200000))
tf = txBox.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = '2005~2024'
p.font.size = Pt(13)
p.font.bold = True
p.font.name = get_font()
p.font.color.rgb = BLACK

add_paragraph(tf, '삼성전자, 삼성SDI, 삼성SDC, CJ, 현대모비스 등', font_size=12, color=GRAY_TEXT, space_before=2)
add_paragraph(tf, '대기업 SI 프로젝트 수행', font_size=12, color=GRAY_TEXT, space_after=12)

add_paragraph(tf, '2025~현재', font_size=13, bold=True, color=BLACK, space_before=8)
add_paragraph(tf, 'asolve(에이솔브) 설립', font_size=12, color=GRAY_TEXT, space_before=2)
add_paragraph(tf, '웹 서비스 기획 · 설계 · 개발 · 운영', font_size=12, color=GRAY_TEXT)

# Right column - Core Skills
right_x = Emu(6400000)

add_textbox(slide, right_x, col_top, Emu(5000000), Emu(300000),
            '핵심 역량', font_size=16, bold=True)

div2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, Emu(1800000), Emu(365760), Emu(27432))
div2.fill.solid()
div2.fill.fore_color.rgb = BLACK
div2.line.fill.background()

skills = [
    '외부 시스템 API 및 인터페이스 연동 전문',
    '데이터 수집 · 관리 · 대시보드 시스템 구축',
    '기획부터 배포까지 전 과정 수행 가능',
    'AI 기반 업무 자동화 솔루션 개발'
]
txBox2 = slide.shapes.add_textbox(right_x, Emu(1900000), Emu(5400000), Emu(1500000))
tf2 = txBox2.text_frame
tf2.word_wrap = True
for i, skill in enumerate(skills):
    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    p.text = '•  ' + skill
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY_TEXT
    p.font.name = get_font()
    p.space_after = Pt(6)

# OCP box
ocp_box = add_rect(slide, right_x, Emu(3600000), Emu(5200000), Emu(350000), fill_color=GRAY_BG)
add_textbox(slide, Emu(6500000), Emu(3620000), Emu(1000000), Emu(300000),
            '자격', font_size=11, color=GRAY_LIGHT)
add_textbox(slide, Emu(7200000), Emu(3620000), Emu(4000000), Emu(300000),
            'OCP (Oracle Certified Professional)', font_size=13, bold=True)

# Homepage
add_textbox(slide, right_x, Emu(4050000), Emu(5200000), Emu(250000),
            '홈페이지:  https://asolv.ai', font_size=12, color=GRAY_LIGHT)


# ═══════════════════════════════════════
# SLIDE 3: Project History
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, '2. 프로젝트 이력')
add_slide_number(slide, 3)
add_title_subtitle(slide, '프로젝트 이력', '21년간 수행한 주요 프로젝트')

# Table
from pptx.util import Cm
rows = 8
cols = 5
table_shape = slide.shapes.add_table(rows, cols, Emu(585216), Emu(1500000), Emu(11200000), Emu(4800000))
table = table_shape.table

# Column widths
table.columns[0].width = Emu(3800000)
table.columns[1].width = Emu(2000000)
table.columns[2].width = Emu(1800000)
table.columns[3].width = Emu(1800000)
table.columns[4].width = Emu(1800000)

headers = ['프로젝트', '기간', '발주처', '역할', '환경']
data = [
    ['삼성전자 NPLM 시스템 구축', '2024~2025', '삼성전자', '인터페이스 리딩', 'Java, Oracle'],
    ['삼성 SDC QMS 2.0 구축', '2023', '삼성 SDC', '개발', 'Java, Vue.js'],
    ['삼성전자 CPCex 사용성 개선', '2022~2023', '삼성전자', '개발', 'Java, Oracle'],
    ['CJ PMS 구축', '2021~2022', 'CJ', '개발', 'Java, Oracle'],
    ['삼성 SDI 표준시스템 고도화', '2020~2021', '삼성 SDI', '개발', 'Java, Oracle'],
    ['현대모비스 연구개발 고도화', '2020', '현대모비스', '개발', 'Java, Spring'],
    ['삼성 SDI PLM 시리즈 (6건)', '2018~2020', '삼성 SDI', '개발', 'Java, Oracle'],
]

for i, h in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = h
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(10)
        paragraph.font.bold = True
        paragraph.font.color.rgb = GRAY_LIGHT
        paragraph.font.name = get_font()
    cell.fill.solid()
    cell.fill.fore_color.rgb = WHITE

for r, row in enumerate(data):
    for c, val in enumerate(row):
        cell = table.cell(r + 1, c)
        cell.text = val
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.name = get_font()
            paragraph.font.color.rgb = BLACK if c == 0 else GRAY_TEXT
            if c == 0:
                paragraph.font.bold = True
            if c == 4:
                paragraph.font.size = Pt(10)
                paragraph.font.color.rgb = GRAY_LIGHT
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE


# ═══════════════════════════════════════
# SLIDE 4: External System Experience
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, '3. 외부 시스템 연계 경험')
add_slide_number(slide, 4)
add_title_subtitle(slide, '외부 시스템 연계 경험', '이번 프로젝트와 직결되는 API · 인터페이스 연동 경험')

# Cards - top row
cards_top = [
    ('삼성전자 NPLM', '2024~2025', ['메모리사업부 150여 개 인터페이스\n직접 리딩 및 수행', '외부 시스템 간 데이터 연계\n아키텍처 설계']),
    ('CJ PMS 구축', '2021~2022', ['프로젝트 관리 시스템 신규 구축', '외부 시스템 연동, 대시보드,\n필터링, 엑셀 내보내기']),
    ('삼성 SDC QMS 2.0', '2023', ['품질관리시스템 내 외부 시스템\n인터페이스 연동']),
]

for i, (title, period, items) in enumerate(cards_top):
    x = Emu(585216 + i * 3750000)
    y = Emu(1400000)
    w = Emu(3550000)
    h = Emu(2000000)
    card = add_rect(slide, x, y, w, h, fill_color=RGBColor(0xFA, 0xFA, 0xFA), border_color=GRAY_BORDER)

    add_textbox(slide, Emu(x + 180000), Emu(y + 150000), Emu(3200000), Emu(250000),
                title, font_size=15, bold=True)
    add_textbox(slide, Emu(x + 180000), Emu(y + 420000), Emu(3200000), Emu(200000),
                period, font_size=10, color=GRAY_LIGHT)

    item_y = y + 700000
    for item in items:
        add_textbox(slide, Emu(x + 180000), Emu(item_y), Emu(3200000), Emu(400000),
                    '•  ' + item, font_size=11, color=GRAY_TEXT)
        item_y += 400000

# Bottom row
cards_bottom = [
    ('삼성 SDI PLM 시리즈', '2018~2020', ['6개 프로젝트 연속 수행', '시스템 간 통합 인터페이스 구축']),
]

x = Emu(585216)
y = Emu(3600000)
w = Emu(5400000)
h = Emu(1700000)
card = add_rect(slide, x, y, w, h, fill_color=RGBColor(0xFA, 0xFA, 0xFA), border_color=GRAY_BORDER)
add_textbox(slide, Emu(x + 180000), Emu(y + 150000), Emu(5000000), Emu(250000),
            '삼성 SDI PLM 시리즈', font_size=15, bold=True)
add_textbox(slide, Emu(x + 180000), Emu(y + 420000), Emu(5000000), Emu(200000),
            '2018~2020', font_size=10, color=GRAY_LIGHT)
add_textbox(slide, Emu(x + 180000), Emu(y + 700000), Emu(5000000), Emu(250000),
            '•  6개 프로젝트 연속 수행', font_size=11, color=GRAY_TEXT)
add_textbox(slide, Emu(x + 180000), Emu(y + 1000000), Emu(5000000), Emu(250000),
            '•  시스템 간 통합 인터페이스 구축', font_size=11, color=GRAY_TEXT)

# Coin arbitrage card (highlighted)
x2 = Emu(6200000)
card2 = add_rect(slide, x2, y, Emu(5585000), h, fill_color=WHITE, border_color=BLACK, border_width=Pt(2))

# Label
label_bg = add_rect(slide, Emu(x2 + 180000), Emu(y + 130000), Emu(1200000), Emu(240000), fill_color=BLACK)
add_textbox(slide, Emu(x2 + 220000), Emu(y + 140000), Emu(1150000), Emu(220000),
            '최근 레퍼런스', font_size=9, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Emu(x2 + 180000), Emu(y + 450000), Emu(5000000), Emu(250000),
            '코인 아비트리지 시스템', font_size=15, bold=True)
add_textbox(slide, Emu(x2 + 180000), Emu(y + 720000), Emu(5000000), Emu(200000),
            '2026.03 · 실제 운영 중', font_size=10, color=GRAY_LIGHT)
add_textbox(slide, Emu(x2 + 180000), Emu(y + 1000000), Emu(5200000), Emu(250000),
            '•  다수의 거래소 API를 웹소켓 기반으로 실시간 연동', font_size=11, color=GRAY_TEXT)
add_textbox(slide, Emu(x2 + 180000), Emu(y + 1280000), Emu(5200000), Emu(250000),
            '•  실시간 가격 비교 · 차익 분석 시스템', font_size=11, color=GRAY_TEXT)


# ═══════════════════════════════════════
# SLIDE 5: Demo
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, '4. 데모 시연')
add_slide_number(slide, 5)
add_title_subtitle(slide, '데모 시연', None)

# Quote box
quote_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(585216), Emu(1100000), Emu(11000000), Emu(400000))
quote_bg.fill.solid()
quote_bg.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
quote_bg.line.fill.background()

quote_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(585216), Emu(1100000), Emu(36576), Emu(400000))
quote_line.fill.solid()
quote_line.fill.fore_color.rgb = BLACK
quote_line.line.fill.background()

add_textbox(slide, Emu(800000), Emu(1150000), Emu(10000000), Emu(300000),
            '"제안서 대신 직접 만들어왔습니다"', font_size=15, bold=True, color=RGBColor(0x33, 0x33, 0x33))

add_textbox(slide, Emu(585216), Emu(1600000), Emu(11000000), Emu(250000),
            '데모 사이트:  https://wish153810.vercel.app', font_size=13, color=BLACK)

# Feature cards (2 rows x 3)
features = [
    ('SAM.gov 연동', 'API Key 입력 시\n실데이터 전환되는 연동 구조'),
    ('대시보드', '실시간 수집 현황, 월별 낙찰\n추이 차트, 마감 D-day'),
    ('입찰 공고', '상태/유형/NAICS/기관 다중 필터\n키워드 검색, 엑셀 다운로드'),
    ('데이터 분석', '경쟁사 낙찰 분석\nNAICS별 시장 분포 차트'),
    ('필터 · 알림', '맞춤 키워드/NAICS 필터\n이메일·SMS·Slack 알림 설정'),
    ('관리자', '사용자 권한 관리\n수집 로그 모니터링'),
]

for i, (title, desc) in enumerate(features):
    col = i % 3
    row = i // 3
    x = Emu(585216 + col * 3750000)
    y = Emu(2050000 + row * 1550000)
    w = Emu(3550000)
    h = Emu(1350000)

    card = add_rect(slide, x, y, w, h, fill_color=RGBColor(0xFA, 0xFA, 0xFA), border_color=GRAY_BORDER)
    add_textbox(slide, Emu(x + 200000), Emu(y + 180000), Emu(3200000), Emu(250000),
                title, font_size=14, bold=True)
    add_textbox(slide, Emu(x + 200000), Emu(y + 520000), Emu(3200000), Emu(600000),
                desc, font_size=11, color=GRAY_LIGHT)

# Bottom bar
bottom_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(585216), Emu(5250000), Emu(11200000), Emu(380000))
bottom_bar.fill.solid()
bottom_bar.fill.fore_color.rgb = BLACK
bottom_bar.line.fill.background()
bottom_bar.adjustments[0] = 0.03
add_textbox(slide, Emu(585216), Emu(5270000), Emu(11200000), Emu(340000),
            'PC / 태블릿 / 모바일 반응형 완전 지원', font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════
# SLIDE 6: Feature List
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, '5. 기능 목록')
add_slide_number(slide, 6)
add_title_subtitle(slide, '기능 목록', '공고 요구사항 기반 전체 기능 정의')

func_items = [
    ('1', '입찰 공고 자동 수집 및 알림', 'SAM.gov API를 통한 주기적 스캔 및 수집 (3시간/6시간/12시간 설정 가능)\n사전공지, 신규, 변경, 마감임박 등 상태별 구분 · 조건 매칭 시 이메일/SMS/Slack 실시간 알림 발송'),
    ('2', '맞춤 필터링 및 선별', 'NAICS 코드, 키워드, 발주기관, 공고유형, 금액범위 다중 필터\n필터 조건에 따른 유망 공고 자동 선별'),
    ('3', '과거 입찰 데이터 분석', '낙찰(Award) 이력 수집 및 경쟁사별 낙찰 건수 · 금액 분석\nNAICS별 시장 분포, 발주 기관별 패턴 분석 차트'),
    ('4', '검색 및 데이터 관리', '국가, 기관, 기간 등 상세 조건 검색 · 수집 데이터 DB 저장 및 이력 관리 · 엑셀(XLSX) 다운로드'),
    ('5', '관리자 기능', '사용자 계정 및 권한 관리 · 필터 키워드·조건 설정 · 알림 채널·주기 환경설정 · 수집 현황 및 시스템 로그 모니터링'),
]

for i, (num, title, desc) in enumerate(func_items):
    y_base = Emu(1400000 + i * 950000)

    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(585216), y_base, Emu(274320), Emu(274320))
    circle.fill.solid()
    circle.fill.fore_color.rgb = BLACK
    circle.line.fill.background()
    add_textbox(slide, Emu(585216), Emu(y_base + 20000), Emu(274320), Emu(250000),
                num, font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Title
    add_textbox(slide, Emu(1050000), y_base, Emu(10000000), Emu(280000),
                title, font_size=14, bold=True)

    # Description
    add_textbox(slide, Emu(1050000), Emu(y_base + 300000), Emu(10500000), Emu(500000),
                desc, font_size=11, color=GRAY_TEXT)


# ═══════════════════════════════════════
# SLIDE 7: Tech Stack
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, '6. 기술 스택')
add_slide_number(slide, 7)
add_title_subtitle(slide, '기술 스택', '프로젝트에 적용할 기술 구성')

tech_sections = [
    ('Frontend', ['Next.js 14 (App Router)', 'TypeScript', 'TailwindCSS', 'Recharts']),
    ('Backend', ['Next.js API Routes', 'Node.js', 'node-cron', 'NextAuth.js']),
    ('외부 연동', ['SAM.gov API v2', 'SendGrid (이메일)', 'CoolSMS (문자)', 'Slack Webhook']),
    ('Database & 인프라', ['PostgreSQL', 'AWS EC2 + RDS', 'Vercel', 'Git']),
]

for i, (section_title, techs) in enumerate(tech_sections):
    col = i % 2
    row = i // 2
    x = Emu(585216 + col * 5700000)
    y = Emu(1400000 + row * 2400000)
    w = Emu(5400000)
    h = Emu(2100000)

    card = add_rect(slide, x, y, w, h, fill_color=RGBColor(0xFA, 0xFA, 0xFA), border_color=GRAY_BORDER)

    add_textbox(slide, Emu(x + 200000), Emu(y + 150000), Emu(4800000), Emu(300000),
                section_title, font_size=16, bold=True)

    # Divider
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(x + 200000), Emu(y + 480000), Emu(365760), Emu(27432))
    div.fill.solid()
    div.fill.fore_color.rgb = BLACK
    div.line.fill.background()

    # Tech labels
    label_y = y + 620000
    label_x = x + 200000
    for tech in techs:
        lbl = add_rect(slide, Emu(label_x), Emu(label_y), Emu(len(tech) * 85000 + 200000), Emu(270000),
                       fill_color=RGBColor(0xF0, 0xF0, 0xF0), border_color=GRAY_BORDER)
        add_textbox(slide, Emu(label_x + 30000), Emu(label_y + 30000), Emu(len(tech) * 85000 + 150000), Emu(210000),
                    tech, font_size=10, bold=True, color=GRAY_SUB, alignment=PP_ALIGN.CENTER)
        label_x += len(tech) * 85000 + 280000
        if label_x > x + 5000000:
            label_x = x + 200000
            label_y += 350000


# ═══════════════════════════════════════
# SLIDE 8: Logical Architecture
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, '7. 시스템 논리 구성도')
add_slide_number(slide, 8)
add_title_subtitle(slide, '시스템 논리 구성도', '전체 시스템 아키텍처 흐름')

center_x = Emu(3100000)
box_w = Emu(6200000)

# User layer
y = Emu(1400000)
user_box = add_rect(slide, center_x, y, box_w, Emu(380000), fill_color=WHITE, border_color=BLACK, border_width=Pt(2))
add_textbox(slide, center_x, Emu(y + 50000), box_w, Emu(280000),
            '사용자 (PC Web - 반응형)', font_size=13, bold=True, alignment=PP_ALIGN.CENTER)

# Arrow
add_textbox(slide, center_x, Emu(1780000), box_w, Emu(300000),
            '↓  HTTPS', font_size=11, color=GRAY_LIGHT, alignment=PP_ALIGN.CENTER)

# Frontend
y2 = Emu(2100000)
fe_box = add_rect(slide, center_x, y2, box_w, Emu(500000), fill_color=BLACK)
add_textbox(slide, center_x, Emu(y2 + 30000), box_w, Emu(200000),
            'Frontend — Next.js', font_size=13, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, center_x, Emu(y2 + 260000), box_w, Emu(200000),
            '대시보드    공고 목록    분석    필터 설정    관리자', font_size=10, color=RGBColor(0x88, 0x88, 0x88), alignment=PP_ALIGN.CENTER)

# Arrow
add_textbox(slide, center_x, Emu(2600000), box_w, Emu(300000),
            '↓  API Routes', font_size=11, color=GRAY_LIGHT, alignment=PP_ALIGN.CENTER)

# Backend
y3 = Emu(2900000)
be_box = add_rect(slide, center_x, y3, box_w, Emu(700000), fill_color=WHITE, border_color=BLACK, border_width=Pt(2))
add_textbox(slide, center_x, Emu(y3 + 40000), box_w, Emu(200000),
            'Backend — Next.js API / Node.js', font_size=13, bold=True, alignment=PP_ALIGN.CENTER)

# Backend sub-items
be_items = ['수집 스케줄러', '필터 매칭 엔진', '알림 발송', '인증/권한']
for i, item in enumerate(be_items):
    ix = Emu(3400000 + i * 1450000)
    iy = Emu(y3 + 350000)
    sub = add_rect(slide, ix, iy, Emu(1300000), Emu(260000), fill_color=GRAY_BG)
    add_textbox(slide, ix, Emu(iy + 30000), Emu(1300000), Emu(200000),
                item, font_size=10, bold=True, color=GRAY_SUB, alignment=PP_ALIGN.CENTER)

# Arrow
add_textbox(slide, center_x, Emu(3600000), box_w, Emu(300000),
            '↓', font_size=14, color=GRAY_LIGHT, alignment=PP_ALIGN.CENTER)

# Bottom 3 boxes
y4 = Emu(3950000)
# SAM.gov
sam_box = add_rect(slide, center_x, y4, Emu(1900000), Emu(600000), fill_color=WHITE, border_color=BLACK, border_width=Pt(2))
add_textbox(slide, center_x, Emu(y4 + 70000), Emu(1900000), Emu(200000),
            'SAM.gov API v2', font_size=12, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, center_x, Emu(y4 + 310000), Emu(1900000), Emu(200000),
            '공고 / 낙찰 데이터', font_size=10, color=GRAY_LIGHT, alignment=PP_ALIGN.CENTER)

# PostgreSQL
pg_box = add_rect(slide, Emu(5200000), y4, Emu(2200000), Emu(600000), fill_color=BLACK)
add_textbox(slide, Emu(5200000), Emu(y4 + 70000), Emu(2200000), Emu(200000),
            'PostgreSQL', font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Emu(5200000), Emu(y4 + 310000), Emu(2200000), Emu(200000),
            '공고 · 낙찰 · 사용자 · 설정 · 로그', font_size=9, color=RGBColor(0x88, 0x88, 0x88), alignment=PP_ALIGN.CENTER)

# Notifications
notif_box = add_rect(slide, Emu(7600000), y4, Emu(1700000), Emu(600000), fill_color=WHITE, border_color=BLACK, border_width=Pt(2))
add_textbox(slide, Emu(7600000), Emu(y4 + 70000), Emu(1700000), Emu(200000),
            '알림 채널', font_size=12, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Emu(7600000), Emu(y4 + 310000), Emu(1700000), Emu(200000),
            'Email / SMS / Slack', font_size=10, color=GRAY_LIGHT, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════
# SLIDE 9: Physical Architecture
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, '8. 시스템 물리 구성도')
add_slide_number(slide, 9)
add_title_subtitle(slide, '시스템 물리 구성도', '사용자 규모에 따른 2가지 인프라 옵션')

# Option A
x_a = Emu(585216)
y_opt = Emu(1400000)
w_half = Emu(5500000)
h_opt = Emu(4200000)

card_a = add_rect(slide, x_a, y_opt, w_half, h_opt, fill_color=WHITE, border_color=BLACK, border_width=Pt(2))
add_textbox(slide, Emu(x_a + 200000), Emu(y_opt + 150000), Emu(3500000), Emu(300000),
            '옵션 A: AWS 구성', font_size=16, bold=True)

badge_a = add_rect(slide, Emu(x_a + 3800000), Emu(y_opt + 170000), Emu(1400000), Emu(240000), fill_color=BLACK)
add_textbox(slide, Emu(x_a + 3800000), Emu(y_opt + 180000), Emu(1400000), Emu(220000),
            '사용자 50명 이상', font_size=9, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

opt_a_items = [
    ('Vercel', 'Frontend 호스팅'),
    ('AWS EC2', 'Backend API + 수집 스케줄러'),
    ('AWS RDS', 'PostgreSQL'),
    ('CloudWatch', '모니터링'),
]

for i, (svc, desc) in enumerate(opt_a_items):
    iy = Emu(y_opt + 600000 + i * 420000)
    svc_box = add_rect(slide, Emu(x_a + 200000), iy, Emu(1500000), Emu(320000), fill_color=GRAY_BG)
    add_textbox(slide, Emu(x_a + 200000), Emu(iy + 50000), Emu(1500000), Emu(220000),
                svc, font_size=11, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Emu(x_a + 1900000), Emu(iy + 60000), Emu(200000), Emu(220000),
                '→', font_size=11, color=GRAY_LIGHT)
    add_textbox(slide, Emu(x_a + 2200000), Emu(iy + 60000), Emu(3000000), Emu(220000),
                desc, font_size=12, color=GRAY_TEXT)

# Cost A
cost_a = add_rect(slide, Emu(x_a + 200000), Emu(y_opt + 3400000), Emu(5100000), Emu(500000), fill_color=GRAY_BG)
add_textbox(slide, Emu(x_a + 200000), Emu(y_opt + 3430000), Emu(5100000), Emu(200000),
            '예상 월 운영비', font_size=10, color=GRAY_LIGHT, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Emu(x_a + 200000), Emu(y_opt + 3650000), Emu(5100000), Emu(250000),
            '10~15만원', font_size=18, bold=True, alignment=PP_ALIGN.CENTER)

# Option B
x_b = Emu(6400000)
card_b = add_rect(slide, x_b, y_opt, w_half, h_opt, fill_color=WHITE, border_color=GRAY_BORDER)
add_textbox(slide, Emu(x_b + 200000), Emu(y_opt + 150000), Emu(3500000), Emu(300000),
            '옵션 B: 경량 구성', font_size=16, bold=True)

badge_b = add_rect(slide, Emu(x_b + 3800000), Emu(y_opt + 170000), Emu(1400000), Emu(240000), fill_color=GRAY_BG)
add_textbox(slide, Emu(x_b + 3800000), Emu(y_opt + 180000), Emu(1400000), Emu(220000),
            '사용자 10명 이내', font_size=9, bold=True, color=GRAY_SUB, alignment=PP_ALIGN.CENTER)

opt_b_items = [
    ('Vercel', 'Frontend + API Routes'),
    ('Vercel Cron', '수집 스케줄러'),
    ('Supabase', 'PostgreSQL'),
    ('Vercel Analytics', '모니터링'),
]

for i, (svc, desc) in enumerate(opt_b_items):
    iy = Emu(y_opt + 600000 + i * 420000)
    svc_box = add_rect(slide, Emu(x_b + 200000), iy, Emu(1500000), Emu(320000), fill_color=GRAY_BG)
    add_textbox(slide, Emu(x_b + 200000), Emu(iy + 50000), Emu(1500000), Emu(220000),
                svc, font_size=11, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Emu(x_b + 1900000), Emu(iy + 60000), Emu(200000), Emu(220000),
                '→', font_size=11, color=GRAY_LIGHT)
    add_textbox(slide, Emu(x_b + 2200000), Emu(iy + 60000), Emu(3000000), Emu(220000),
                desc, font_size=12, color=GRAY_TEXT)

cost_b = add_rect(slide, Emu(x_b + 200000), Emu(y_opt + 3400000), Emu(5100000), Emu(500000), fill_color=GRAY_BG)
add_textbox(slide, Emu(x_b + 200000), Emu(y_opt + 3430000), Emu(5100000), Emu(200000),
            '예상 월 운영비', font_size=10, color=GRAY_LIGHT, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Emu(x_b + 200000), Emu(y_opt + 3650000), Emu(5100000), Emu(250000),
            '3~5만원', font_size=18, bold=True, alignment=PP_ALIGN.CENTER)

# Bottom note
add_textbox(slide, Emu(585216), Emu(5800000), Emu(11200000), Emu(300000),
            '사용자 규모와 예산에 따라 협의하여 결정', font_size=13, color=GRAY_SUB, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════
# SLIDE 10: ERD
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, '9. 데이터 구성도 (ERD)')
add_slide_number(slide, 10)
add_title_subtitle(slide, '데이터 구성도 (ERD)', '현재 기능 기준 데이터베이스 테이블 구성')

erd_tables = [
    ('opportunities', '입찰 공고', ['id  PK', 'notice_id', 'title', 'type, status', 'posted_date', 'response_deadline', 'naics_code', 'agency, office', 'description', 'award_amount', 'awardee', 'sam_url']),
    ('users', '사용자', ['id  PK', 'email', 'name', 'role', 'password_hash', 'created_at']),
    ('filters', '필터 설정', ['id  PK', 'user_id  FK', 'name', 'keywords', 'naics_codes', 'agencies', 'is_active']),
    ('notifications', '알림 이력', ['id  PK', 'user_id  FK', 'opportunity_id  FK', 'channel', 'sent_at', 'status']),
]

erd_tables2 = [
    ('notification_settings', '알림 설정', ['id  PK', 'user_id  FK', 'email_enabled', 'sms_enabled', 'slack_enabled', 'email_address', 'phone_number', 'slack_webhook_url']),
    ('collection_logs', '수집 로그', ['id  PK', 'started_at', 'completed_at', 'status', 'total_fetched', 'new_count', 'error_message']),
    ('bookmarks', '공고 저장', ['id  PK', 'user_id  FK', 'opportunity_id  FK', 'created_at']),
]

# Row 1
for i, (name, label, fields) in enumerate(erd_tables):
    x = Emu(585216 + i * 2900000)
    y = Emu(1400000)
    w = Emu(2700000)

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(x), y, w, Emu(300000))
    header.fill.solid()
    header.fill.fore_color.rgb = BLACK
    header.line.fill.background()
    header.adjustments[0] = 0.04

    add_textbox(slide, Emu(x + 100000), Emu(y + 50000), Emu(2500000), Emu(200000),
                f'{name}  {label}', font_size=10, bold=True, color=WHITE)

    # Fields
    for j, field in enumerate(fields):
        fy = Emu(y + 340000 + j * 210000)
        field_color = GRAY_TEXT
        if 'PK' in field:
            field_color = BLACK
        elif 'FK' in field:
            field_color = BLUE_FK
        add_textbox(slide, Emu(x + 100000), fy, Emu(2500000), Emu(200000),
                    field, font_size=9, color=field_color,
                    font_name='Consolas' if True else get_font())
        # Separator line
        if j < len(fields) - 1:
            sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(x + 50000), Emu(fy + 200000), Emu(w - 100000), Emu(9144))
            sep.fill.solid()
            sep.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
            sep.line.fill.background()

# Row 2
for i, (name, label, fields) in enumerate(erd_tables2):
    x = Emu(585216 + i * 3900000)
    y = Emu(4500000)
    w = Emu(3700000)

    header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(x), y, w, Emu(300000))
    header.fill.solid()
    header.fill.fore_color.rgb = BLACK
    header.line.fill.background()
    header.adjustments[0] = 0.03

    add_textbox(slide, Emu(x + 100000), Emu(y + 50000), Emu(3500000), Emu(200000),
                f'{name}  {label}', font_size=10, bold=True, color=WHITE)

    for j, field in enumerate(fields):
        fy = Emu(y + 340000 + j * 210000)
        field_color = GRAY_TEXT
        if 'PK' in field:
            field_color = BLACK
        elif 'FK' in field:
            field_color = BLUE_FK
        add_textbox(slide, Emu(x + 100000), fy, Emu(3500000), Emu(200000),
                    field, font_size=9, color=field_color,
                    font_name='Consolas')
        if j < len(fields) - 1:
            sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(x + 50000), Emu(fy + 200000), Emu(w - 100000), Emu(9144))
            sep.fill.solid()
            sep.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
            sep.line.fill.background()


# ═══════════════════════════════════════
# SLIDE 11: Schedule
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, '10. 프로젝트 일정')
add_slide_number(slide, 11)
add_title_subtitle(slide, '프로젝트 일정', '30일 기준 주차별 일정 및 액션아이템')

weeks = [
    ('WEEK 1', '기획 · 설계', [
        ('클라이언트', '요구사항 최종 확인, 피드백'),
        ('개발자', '요구사항 정의서 작성'),
        ('개발자', '기능명세서 작성'),
        ('개발자', '화면설계서 작성'),
        ('개발자', 'DB 스키마 설계'),
    ], '요구사항 정의서\n화면 설계서'),
    ('WEEK 2', '백엔드 핵심', [
        ('클라이언트', '화면 설계서 검수·승인'),
        ('클라이언트', 'SAM.gov API Key 발급'),
        ('개발자', 'SAM.gov API 연동'),
        ('개발자', 'DB 구축'),
        ('개발자', '수집 스케줄러, 필터 엔진'),
    ], 'API 연동 완료\n개발서버 오픈'),
    ('WEEK 3', '프론트 + 알림', [
        ('클라이언트', '개발서버 확인·피드백'),
        ('개발자', '전체 페이지 개발'),
        ('개발자', '알림 연동'),
        ('개발자', '엑셀 내보내기, 분석 차트'),
    ], '전 기능 개발 완료'),
    ('WEEK 4', '통합·검수·납품', [
        ('클라이언트', '최종 검수, 피드백'),
        ('개발자', 'QA, 버그 수정'),
        ('개발자', '서버 배포'),
        ('개발자', '산출물 정리'),
    ], '소스코드, 기획서\n일체 납품'),
]

for i, (week, title, items, output) in enumerate(weeks):
    x = Emu(585216 + i * 2950000)
    y_base = Emu(1400000)

    # Timeline bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(x), y_base, Emu(27432), Emu(4200000))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLACK
    bar.line.fill.background()

    add_textbox(slide, Emu(x + 100000), y_base, Emu(2500000), Emu(200000),
                week, font_size=9, bold=True, color=BLACK)
    add_textbox(slide, Emu(x + 100000), Emu(y_base + 220000), Emu(2500000), Emu(250000),
                title, font_size=14, bold=True)

    for j, (role, task) in enumerate(items):
        iy = Emu(y_base + 600000 + j * 280000)
        dot_color = BLUE_FK if role == '클라이언트' else BLACK
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(x + 120000), Emu(iy + 50000), Emu(54864), Emu(54864))
        dot.fill.solid()
        dot.fill.fore_color.rgb = dot_color
        dot.line.fill.background()
        add_textbox(slide, Emu(x + 230000), iy, Emu(2600000), Emu(250000),
                    task, font_size=10, color=GRAY_TEXT)

    # Output box
    out_y = Emu(y_base + 3200000)
    out_box = add_rect(slide, Emu(x + 80000), out_y, Emu(2700000), Emu(600000), fill_color=GRAY_BG)
    add_textbox(slide, Emu(x + 150000), Emu(out_y + 60000), Emu(2600000), Emu(500000),
                output, font_size=9, bold=True, color=GRAY_SUB)

# Legend
add_textbox(slide, Emu(585216), Emu(5900000), Emu(3000000), Emu(200000),
            '●  클라이언트 액션      ●  개발자 액션', font_size=9, color=GRAY_LIGHT)


# ═══════════════════════════════════════
# SLIDE 12: Discussion Points
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, '11. 협의 사항')
add_slide_number(slide, 12)
add_title_subtitle(slide, '협의 사항', '사전 확인이 필요한 항목')

discussion = [
    ('1', 'SAM.gov API Key 발급', '클라이언트 명의로 미국 정부 사이트에서 발급 필요\n발급 절차 가이드 지원 가능'),
    ('2', '알림 채널 확정', '이메일 / SMS / Slack 중 실제 필요한 채널\nSMS, 카카오는 별도 서비스 가입 및 발송 비용 발생'),
    ('3', '서버 · 호스팅 비용', '운영 서버 및 DB 비용은 클라이언트 카드 결제 구조\n월 예상: 3~15만원 (옵션에 따라 상이)'),
    ('4', '사용자 규모', '내부 직원 전용인지, 외부 고객 서비스인지\n아키텍처 옵션 선택 기준'),
    ('5', '디자인 산출물', '화면 설계서 (이미지/PDF) 형태로 납품'),
    ('6', '유지보수', '납품 후 1개월 무상 하자보수 제안\n이후 유지보수는 별도 협의'),
]

for i, (num, title, desc) in enumerate(discussion):
    col = i // 3
    row = i % 3
    x = Emu(585216 + col * 5800000)
    y = Emu(1400000 + row * 1500000)

    # Number box
    num_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(x), y, Emu(274320), Emu(274320))
    num_box.fill.background()
    num_box.line.color.rgb = BLACK
    num_box.line.width = Pt(2)
    num_box.adjustments[0] = 0.1
    add_textbox(slide, Emu(x), Emu(y + 20000), Emu(274320), Emu(250000),
                num, font_size=11, bold=True, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Emu(x + 400000), y, Emu(5000000), Emu(280000),
                title, font_size=14, bold=True)
    add_textbox(slide, Emu(x + 400000), Emu(y + 320000), Emu(5000000), Emu(800000),
                desc, font_size=11, color=GRAY_LIGHT)

    # Separator
    if row < 2:
        sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(x), Emu(y + 1300000), Emu(5200000), Emu(9144))
        sep.fill.solid()
        sep.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
        sep.line.fill.background()


# ═══════════════════════════════════════
# SLIDE 13: Closing
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, '')
add_slide_number(slide, 13)

# Center line
cline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(5850000), Emu(2000000), Emu(548640), Emu(36576))
cline.fill.solid()
cline.fill.fore_color.rgb = BLACK
cline.line.fill.background()

add_textbox(slide, Emu(0), Emu(2200000), prs.slide_width, Emu(500000),
            '감사합니다', font_size=36, bold=True, alignment=PP_ALIGN.CENTER)

# Info boxes
add_textbox(slide, Emu(3500000), Emu(3200000), Emu(2500000), Emu(180000),
            '데모 사이트', font_size=10, color=GRAY_LIGHT, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Emu(3500000), Emu(3400000), Emu(2500000), Emu(250000),
            'wish153810.vercel.app', font_size=14, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Emu(6500000), Emu(3200000), Emu(2500000), Emu(180000),
            '홈페이지', font_size=10, color=GRAY_LIGHT, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Emu(6500000), Emu(3400000), Emu(2500000), Emu(250000),
            'asolv.ai', font_size=14, bold=True, alignment=PP_ALIGN.CENTER)

# Name card
name_card = add_rect(slide, Emu(4500000), Emu(4200000), Emu(3400000), Emu(800000), fill_color=RGBColor(0xFA, 0xFA, 0xFA), border_color=GRAY_BORDER)
add_textbox(slide, Emu(4500000), Emu(4300000), Emu(3400000), Emu(300000),
            'asolve (에이솔브)', font_size=16, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Emu(4500000), Emu(4600000), Emu(3400000), Emu(250000),
            '대표 · 개발자 유영주', font_size=12, color=GRAY_LIGHT, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════
# SAVE
# ═══════════════════════════════════════
output_path = os.path.join(os.path.dirname(__file__), 'USFK-BidTrack-Proposal.pptx')
prs.save(output_path)
print(f'PPTX saved: {output_path}')
